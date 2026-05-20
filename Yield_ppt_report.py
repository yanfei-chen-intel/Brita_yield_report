"""
Brita Yield Report – PowerPoint generator
"""

import io
import json
import pathlib
import time
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# Config
# ---------------------------------------------------------------------------
CONFIG_FILE = pathlib.Path("input") / "config.json"

def _load_config() -> dict:
    if CONFIG_FILE.exists():
        with open(CONFIG_FILE, "r") as f:
            return json.load(f)
    return {}


def _wait_for_files(files: list, check_period: int, check_count):
    """Block until all *files* exist, polling every *check_period* seconds.

    check_count: max retries (None = unlimited).
    """
    missing = [p for p in files if not pathlib.Path(p).exists()]
    if not missing:
        return

    attempt = 0
    while missing:
        names = [pathlib.Path(p).name for p in missing]
        print(f"  [wait] Missing: {names}  – retrying in {check_period}s "
              f"(attempt {attempt + 1}"
              + (f"/{check_count})" if check_count else ")"))
        time.sleep(check_period)
        attempt += 1
        missing = [p for p in missing if not pathlib.Path(p).exists()]
        if check_count is not None and attempt >= check_count:
            raise TimeoutError(
                f"Files still missing after {attempt} attempts: "
                + ", ".join(str(p) for p in missing)
            )
    print(f"  [wait] All files ready.")


# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
OUTPUT_DIR         = pathlib.Path("output")
XEUS_CSV           = OUTPUT_DIR / "xeus_test_result.csv"
LOT_WAFER_RESULT   = OUTPUT_DIR / "lot_wafer_search_result.csv"
SOC_SCAN_HRY_CSV    = OUTPUT_DIR / "SOC_SCAN" / "HRY_soc_scan.csv"
SOC_SCAN_HRY_XLSX   = OUTPUT_DIR / "SOC_SCAN" / "HRY_soc_scan.xlsx"
SOC_SCAN_YIELD_HTML = OUTPUT_DIR / "SOC_SCAN" / "yield_plot.html"
MBIST_AGG_CSV       = OUTPUT_DIR / "MBIST" / "MBIST_units_batch_summary_agg.csv"
MBIST_FILTERED_CSV  = OUTPUT_DIR / "MBIST" / "MBIST_units_batch_summary_filtered.csv"
ARRGT_SPLIT1_JRN    = OUTPUT_DIR / "ARRGT" / "ARRGT_Split1_Bivariate.jrn"
ARRGT_SPLIT2_JRN    = OUTPUT_DIR / "ARRGT" / "ARRGT_Split2_Bivariate.jrn"
ARRGT_SPLIT1_PPTX   = OUTPUT_DIR / "ARRGT" / "ARRGT_Split1_Bivariate.pptx"
ARRGT_SPLIT2_PPTX   = OUTPUT_DIR / "ARRGT" / "ARRGT_Split2_Bivariate.pptx"
ARRGT_SPLIT1_PNG    = OUTPUT_DIR / "ARRGT" / "ARRGT_Split1_Bivariate_slide1.png"
ARRGT_SPLIT2_PNG    = OUTPUT_DIR / "ARRGT" / "ARRGT_Split2_Bivariate_slide1.png"
MIO_FILTERED_CSV    = OUTPUT_DIR / "MIO"   / "MIO_filtered_yield.csv"
MIO_YIELD_CSV       = OUTPUT_DIR / "MIO"   / "MIO_yield.csv"
PPT_OUT             = OUTPUT_DIR / "Brita_yield_report.pptx"

# ---------------------------------------------------------------------------
# Colour palette (Intel-ish)
# ---------------------------------------------------------------------------
INTEL_BLUE  = RGBColor(0x00, 0x68, 0xB5)
LIGHT_BLUE  = RGBColor(0xD6, 0xEA, 0xF8)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GREY   = RGBColor(0x2C, 0x2C, 0x2C)
MID_GREY    = RGBColor(0x60, 0x60, 0x60)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def _add_textbox(slide, left, top, width, height, text,
                 font_size=12, bold=False, color=DARK_GREY,
                 align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txb


def _add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _add_table(slide, rows, cols, left, top, width, height):
    return slide.shapes.add_table(rows, cols, left, top, width, height).table


def _style_cell(cell, text, font_size=11, bold=False,
                fg=DARK_GREY, bg=None, align=PP_ALIGN.LEFT):
    cell.text = text
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = fg


# ---------------------------------------------------------------------------
# Page 1 – Test Run Summary
# ---------------------------------------------------------------------------
def build_test_run_summary(prs: Presentation):
    slide_layout = prs.slide_layouts[6]   # blank
    slide = prs.slides.add_slide(slide_layout)

    # ── Header bar ──────────────────────────────────────────────────────────
    _add_rect(slide, 0, 0, SLIDE_W, Inches(1.0), INTEL_BLUE)
    _add_textbox(
        slide,
        Inches(0.3), Inches(0.15),
        Inches(12), Inches(0.7),
        "Brita Yield Report  |  Test Run Summary",
        font_size=24, bold=True, color=WHITE,
    )

    # ── Load data ────────────────────────────────────────────────────────────
    df_xeus = pd.read_csv(XEUS_CSV, dtype=str)
    df_xeus.columns = df_xeus.columns.str.strip()
    row = df_xeus.iloc[0]

    lot        = row.get("Lot",        "N/A")
    operation  = row.get("Operation",  "N/A")
    part_type  = row.get("PartType",   "N/A")
    stpl_path  = row.get("StplPath",   "N/A")

    df_lw = pd.read_csv(LOT_WAFER_RESULT, dtype=str)
    df_lw.columns = df_lw.columns.str.strip()
    wafer_ids = sorted(df_lw["WAFER_ID"].dropna().unique().tolist(),
                       key=lambda x: int(x) if x.isdigit() else x)

    # ── Info table (left panel) ──────────────────────────────────────────────
    fields = [
        ("Lot",        lot),
        ("Operation",  operation),
        ("Part Type",  part_type),
        ("STPL Path",  stpl_path),
    ]

    tbl_left   = Inches(0.4)
    tbl_top    = Inches(1.25)
    tbl_width  = Inches(8.5)
    tbl_height = Inches(0.45 * len(fields))

    table = _add_table(slide, len(fields), 2,
                       tbl_left, tbl_top, tbl_width, tbl_height)

    table.columns[0].width = Inches(1.8)
    table.columns[1].width = Inches(6.7)

    for i, (label, value) in enumerate(fields):
        bg = LIGHT_BLUE if i % 2 == 0 else WHITE
        _style_cell(table.cell(i, 0), label,
                    font_size=11, bold=True, fg=INTEL_BLUE, bg=bg)
        _style_cell(table.cell(i, 1), value,
                    font_size=11, bold=False, fg=DARK_GREY, bg=bg)

    # ── Wafer IDs panel (right) ──────────────────────────────────────────────
    panel_left = Inches(9.2)
    panel_top  = Inches(1.25)
    panel_w    = Inches(3.8)
    panel_h    = Inches(5.8)

    _add_rect(slide, panel_left, panel_top, panel_w, panel_h, LIGHT_BLUE)

    _add_textbox(
        slide,
        panel_left + Inches(0.15), panel_top + Inches(0.1),
        panel_w - Inches(0.3), Inches(0.4),
        f"Wafer IDs  ({len(wafer_ids)} wafers)",
        font_size=12, bold=True, color=INTEL_BLUE,
    )

    wafer_text = "\n".join(wafer_ids)
    _add_textbox(
        slide,
        panel_left + Inches(0.15), panel_top + Inches(0.55),
        panel_w - Inches(0.3), panel_h - Inches(0.7),
        wafer_text,
        font_size=11, color=DARK_GREY,
    )

    # ── Footer ───────────────────────────────────────────────────────────────
    _add_rect(slide,
              0, SLIDE_H - Inches(0.3),
              SLIDE_W, Inches(0.3), INTEL_BLUE)
    _add_textbox(
        slide,
        Inches(0.3), SLIDE_H - Inches(0.28),
        Inches(12), Inches(0.25),
        "Intel Confidential",
        font_size=8, color=WHITE,
    )


# ---------------------------------------------------------------------------
# Page 2 – SOC_SCAN Results
# ---------------------------------------------------------------------------
def _build_yield_chart(df: pd.DataFrame) -> io.BytesIO:
    """Recreate the yield box/bar chart from HRY_soc_scan.csv as a PNG."""
    df = df.copy()
    df["Yield_pct"] = df["Yield"].str.rstrip("%").astype(float)
    df_sorted = df.sort_values("Yield_pct", ascending=True)

    fig, ax = plt.subplots(figsize=(4.8, 6.5))
    colors = ["#d32f2f" if y < 99 else "#0068B5" for y in df_sorted["Yield_pct"]]
    bars = ax.barh(df_sorted["Indicator"], df_sorted["Yield_pct"],
                   color=colors, height=0.7, edgecolor="white", linewidth=0.4)

    ax.set_xlabel("Yield (%)", fontsize=8)
    ax.set_title("Yield by Indicator", fontsize=10, fontweight="bold", color="#2C2C2C")
    ax.set_xlim(0, 105)
    ax.tick_params(axis="y", labelsize=6)
    ax.tick_params(axis="x", labelsize=7)
    ax.xaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: f"{x:.0f}%"))
    ax.axvline(99, color="orange", linestyle="--", linewidth=0.8, label="99% threshold")
    ax.grid(axis="x", linestyle="--", linewidth=0.4, alpha=0.6)
    ax.set_facecolor("#F5F8FC")
    fig.patch.set_facecolor("white")

    # Color y-axis tick labels red for yield < 99%
    fig.canvas.draw()
    yield_vals = df_sorted["Yield_pct"].values
    for tick, val in zip(ax.get_yticklabels(), yield_vals):
        if val < 99:
            tick.set_color("#d32f2f")

    pass_patch = mpatches.Patch(color="#0068B5", label="≥ 99%")
    fail_patch = mpatches.Patch(color="#d32f2f", label="< 99%")
    ax.legend(handles=[pass_patch, fail_patch], fontsize=7, loc="lower right")

    plt.tight_layout(pad=0.5)
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
    buf.seek(0)
    plt.close(fig)
    return buf


def build_SOC_SCAN_result(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ── Header ───────────────────────────────────────────────────────────────
    _add_rect(slide, 0, 0, SLIDE_W, Inches(1.0), INTEL_BLUE)
    _add_textbox(
        slide,
        Inches(0.3), Inches(0.15),
        Inches(12), Inches(0.7),
        "Brita Yield Report  |  SOC_SCAN Results",
        font_size=24, bold=True, color=WHITE,
    )

    # ── Footer ───────────────────────────────────────────────────────────────
    _add_rect(slide, 0, SLIDE_H - Inches(0.3), SLIDE_W, Inches(0.3), INTEL_BLUE)
    _add_textbox(
        slide,
        Inches(0.3), SLIDE_H - Inches(0.28),
        Inches(12), Inches(0.25),
        "Intel Confidential",
        font_size=8, color=WHITE,
    )

    # ── Load data ────────────────────────────────────────────────────────────
    df = pd.read_csv(SOC_SCAN_HRY_CSV, dtype=str)
    df.columns = df.columns.str.strip()

    content_top = Inches(1.08)
    label_h     = Inches(0.30)
    chart_top   = content_top + label_h
    chart_h     = SLIDE_H - chart_top - Inches(0.38)

    # ── Chart section label (left) ───────────────────────────────────────────
    _add_textbox(
        slide,
        Inches(0.2), content_top, Inches(8.4), label_h,
        "Yield Distribution by Indicator",
        font_size=10, bold=True, color=INTEL_BLUE,
    )

    # ── Yield chart (left ~64%) ───────────────────────────────────────────────
    chart_buf = _build_yield_chart(df)
    slide.shapes.add_picture(
        chart_buf,
        Inches(0.2), chart_top,
        Inches(8.4), chart_h,
    )

    # ── Attachment panel (right ~34%) ────────────────────────────────────────
    attach_left  = Inches(8.8)
    attach_top   = content_top
    attach_w     = Inches(4.3)
    attach_h     = chart_h + label_h

    _add_rect(slide, attach_left, attach_top, attach_w, attach_h, LIGHT_BLUE)

    _add_textbox(
        slide,
        attach_left + Inches(0.15), attach_top + Inches(0.12),
        attach_w - Inches(0.3), Inches(0.30),
        "HRY SOC_SCAN Detail Data",
        font_size=10, bold=True, color=INTEL_BLUE,
    )
    _add_textbox(
        slide,
        attach_left + Inches(0.15), attach_top + Inches(0.48),
        attach_w - Inches(0.3), Inches(0.22),
        "File: HRY_soc_scan.csv",
        font_size=9, bold=False, color=MID_GREY,
    )
    _add_textbox(
        slide,
        attach_left + Inches(0.15), attach_top + Inches(0.72),
        attach_w - Inches(0.3), Inches(0.22),
        "Double-click the icon below to open the data",
        font_size=8, bold=False, color=MID_GREY,
    )


def _embed_csv_ole(pptx_path: pathlib.Path, csv_path: pathlib.Path,
                   slide_index: int = 1,
                   left_in: float = 8.8, top_in: float = 2.8,
                   width_in: float = 2.2, height_in: float = 1.6):
    """Use PowerPoint COM to embed csv_path as an OLE icon on the given slide."""
    try:
        import win32com.client
    except ImportError:
        print("pywin32 not available – skipping OLE embedding.")
        return

    pptx_abs = str(pptx_path.resolve())
    csv_abs  = str(csv_path.resolve())

    ppt_app = win32com.client.Dispatch("PowerPoint.Application")
    ppt_app.Visible = True
    try:
        prs_com = ppt_app.Presentations.Open(pptx_abs, ReadOnly=False,
                                              Untitled=False, WithWindow=False)
        slide = prs_com.Slides(slide_index + 1)  # COM is 1-indexed

        slide.Shapes.AddOLEObject(
            Left=left_in * 72, Top=top_in * 72,
            Width=width_in * 72, Height=height_in * 72,
            FileName=csv_abs,
            DisplayAsIcon=True,
            IconLabel=csv_path.name,
        )
        prs_com.Save()
        prs_com.Close()
        print(f"OLE attachment embedded: {csv_path.name} on slide {slide_index + 1}")
    except Exception as exc:
        print(f"OLE embedding failed ({csv_path.name}): {exc}")
    finally:
        try:
            ppt_app.Quit()
        except Exception:
            pass


# ---------------------------------------------------------------------------
# Page 3 – MBIST Results
# ---------------------------------------------------------------------------
def _shorten_indicator(name: str) -> str:
    """Strip common SOCARRAY:: prefix and ::Y suffix for compact display."""
    name = name.replace("SOCARRAY::", "")
    if name.endswith("::Y"):
        name = name[:-3]
    return name


def _build_mbist_chart(df_agg: pd.DataFrame) -> io.BytesIO:
    """Horizontal bar chart of Avg_%Impact by indicator from the agg CSV."""
    df = df_agg.copy()
    df["Short"] = df["Indicator"].apply(_shorten_indicator)
    df["Avg"]   = pd.to_numeric(df["Avg_%Impact"], errors="coerce").fillna(0)
    df_sorted   = df.sort_values("Avg", ascending=True)

    fig, ax = plt.subplots(figsize=(5.2, 3.8))
    colors = ["#d32f2f" if v < 99 else "#0068B5" for v in df_sorted["Avg"]]
    ax.barh(df_sorted["Short"], df_sorted["Avg"],
            color=colors, height=0.55, edgecolor="white", linewidth=0.4)

    ax.set_xlabel("Avg %Impact on Tested Die (%)", fontsize=8)
    ax.set_title("Avg %Impact by Indicator (all wafers)", fontsize=9,
                 fontweight="bold", color="#2C2C2C")
    ax.set_xlim(0, 105)
    ax.tick_params(axis="y", labelsize=7)
    ax.tick_params(axis="x", labelsize=7)
    ax.xaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: f"{x:.0f}%"))
    ax.axvline(99, color="orange", linestyle="--", linewidth=0.8)
    ax.grid(axis="x", linestyle="--", linewidth=0.4, alpha=0.6)
    ax.set_facecolor("#F5F8FC")
    fig.patch.set_facecolor("white")

    pass_patch = mpatches.Patch(color="#0068B5", label="≥ 99%")
    fail_patch = mpatches.Patch(color="#d32f2f", label="< 99%")
    ax.legend(handles=[pass_patch, fail_patch], fontsize=7, loc="lower right")

    plt.tight_layout(pad=0.5)
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
    buf.seek(0)
    plt.close(fig)
    return buf


def build_MBIST_result(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ── Header ───────────────────────────────────────────────────────────────
    _add_rect(slide, 0, 0, SLIDE_W, Inches(1.0), INTEL_BLUE)
    _add_textbox(
        slide,
        Inches(0.3), Inches(0.15),
        Inches(12), Inches(0.7),
        "Brita Yield Report  |  MBIST Results",
        font_size=24, bold=True, color=WHITE,
    )

    # ── Footer ───────────────────────────────────────────────────────────────
    _add_rect(slide, 0, SLIDE_H - Inches(0.3), SLIDE_W, Inches(0.3), INTEL_BLUE)
    _add_textbox(
        slide,
        Inches(0.3), SLIDE_H - Inches(0.28),
        Inches(12), Inches(0.25),
        "Intel Confidential",
        font_size=8, color=WHITE,
    )

    # ── Load agg data ────────────────────────────────────────────────────────
    df_agg = pd.read_csv(MBIST_AGG_CSV, dtype=str)
    df_agg.columns = df_agg.columns.str.strip()

    content_top = Inches(1.08)
    label_h     = Inches(0.30)
    left_w      = Inches(8.4)
    right_left  = Inches(8.8)
    right_w     = Inches(4.3)
    content_h   = SLIDE_H - content_top - Inches(0.38)

    # ── Left: agg summary table ───────────────────────────────────────────────
    _add_textbox(
        slide,
        Inches(0.2), content_top, left_w, label_h,
        "MBIST Aggregated Summary  (MBIST_units_batch_summary_agg.csv)",
        font_size=10, bold=True, color=INTEL_BLUE,
    )

    tbl_top  = content_top + label_h
    tbl_h    = Inches(0.38 * (len(df_agg) + 1))   # header + data rows
    tbl_cols = ["Indicator", "Wafer_Count", "Sum_%Impact", "Avg_%Impact"]
    col_w    = [Inches(4.5), Inches(0.9), Inches(1.5), Inches(1.5)]

    tbl = _add_table(slide, len(df_agg) + 1, len(tbl_cols),
                     Inches(0.2), tbl_top, left_w, tbl_h)
    for j, (col, w) in enumerate(zip(tbl_cols, col_w)):
        tbl.columns[j].width = w
        _style_cell(tbl.cell(0, j), col,
                    font_size=9, bold=True, fg=WHITE, bg=INTEL_BLUE,
                    align=PP_ALIGN.CENTER)

    for i, row in df_agg.iterrows():
        bg = LIGHT_BLUE if i % 2 == 0 else WHITE
        for j, col in enumerate(tbl_cols):
            val = str(row[col])
            if col in ("Sum_%Impact", "Avg_%Impact"):
                try:
                    val = f"{float(val):.2f}%"
                except ValueError:
                    pass
            elif col == "Indicator":
                val = _shorten_indicator(val)
            _style_cell(tbl.cell(i + 1, j), val,
                        font_size=9, fg=DARK_GREY, bg=bg)

    # ── Left: bar chart below table ───────────────────────────────────────────
    chart_top = tbl_top + tbl_h + Inches(0.15)
    chart_h   = SLIDE_H - chart_top - Inches(0.38)

    _add_textbox(
        slide,
        Inches(0.2), chart_top - Inches(0.28), left_w, label_h,
        "Avg %Impact by Indicator",
        font_size=10, bold=True, color=INTEL_BLUE,
    )

    chart_buf = _build_mbist_chart(df_agg)
    slide.shapes.add_picture(
        chart_buf,
        Inches(0.2), chart_top,
        left_w, chart_h,
    )

    # ── Right: attachment panel ───────────────────────────────────────────────
    _add_rect(slide, right_left, content_top, right_w, content_h, LIGHT_BLUE)

    _add_textbox(
        slide,
        right_left + Inches(0.15), content_top + Inches(0.12),
        right_w - Inches(0.3), Inches(0.30),
        "MBIST Detail Data",
        font_size=10, bold=True, color=INTEL_BLUE,
    )
    _add_textbox(
        slide,
        right_left + Inches(0.15), content_top + Inches(0.48),
        right_w - Inches(0.3), Inches(0.22),
        "File: MBIST_units_batch_summary_filtered.csv",
        font_size=8, bold=False, color=MID_GREY,
    )
    _add_textbox(
        slide,
        right_left + Inches(0.15), content_top + Inches(0.72),
        right_w - Inches(0.3), Inches(0.22),
        "Double-click the icon below to open the data",
        font_size=8, bold=False, color=MID_GREY,
    )


def _export_slide_as_image(src_pptx: pathlib.Path, out_png: pathlib.Path,
                            width_px: int = 1600, height_px: int = 900) -> bool:
    """Export the first slide of *src_pptx* to *out_png* via PowerPoint COM."""
    try:
        import win32com.client
    except ImportError:
        print("pywin32 not available – skipping slide image export.")
        return False
    if not src_pptx.exists():
        print(f"  [warn] Source PPTX not found: {src_pptx}")
        return False
    abs_src = str(src_pptx.resolve())
    abs_out = str(out_png.resolve())
    app = win32com.client.Dispatch("PowerPoint.Application")
    app.Visible = True
    try:
        pres = app.Presentations.Open(abs_src, ReadOnly=True,
                                       Untitled=False, WithWindow=False)
        pres.Slides(1).Export(abs_out, "PNG", width_px, height_px)
        pres.Close()
        print(f"  Slide exported → {out_png.name}")
        return True
    except Exception as exc:
        print(f"  [warn] Slide export failed: {exc}")
        return False


# ---------------------------------------------------------------------------
# Page 4 – ARRGT Results
# ---------------------------------------------------------------------------
def build_ARRGT_result(prs: Presentation,
                        img1: pathlib.Path = None,
                        img2: pathlib.Path = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ── Header ───────────────────────────────────────────────────────────────
    _add_rect(slide, 0, 0, SLIDE_W, Inches(1.0), INTEL_BLUE)
    _add_textbox(
        slide,
        Inches(0.3), Inches(0.15),
        Inches(12), Inches(0.7),
        "Brita Yield Report  |  ARRGT Results",
        font_size=24, bold=True, color=WHITE,
    )

    # ── Footer ───────────────────────────────────────────────────────────────
    _add_rect(slide, 0, SLIDE_H - Inches(0.3), SLIDE_W, Inches(0.3), INTEL_BLUE)
    _add_textbox(
        slide,
        Inches(0.3), SLIDE_H - Inches(0.28),
        Inches(12), Inches(0.25),
        "Intel Confidential",
        font_size=8, color=WHITE,
    )

    content_top = Inches(1.08)
    content_h   = SLIDE_H - content_top - Inches(0.38)
    panel_w     = Inches(6.3)
    gap         = Inches(0.23)
    left1       = Inches(0.2)
    left2       = left1 + panel_w + gap
    img_top     = content_top + Inches(1.1)
    img_h       = Inches(3.2)
    img_w       = panel_w - Inches(0.3)

    # ── Split1 Bivariate panel ───────────────────────────────────────────────
    _add_rect(slide, left1, content_top, panel_w, content_h, LIGHT_BLUE)
    _add_textbox(
        slide,
        left1 + Inches(0.15), content_top + Inches(0.12),
        panel_w - Inches(0.3), Inches(0.30),
        "ARRGT Split1 Bivariate Analysis",
        font_size=10, bold=True, color=INTEL_BLUE,
    )
    _add_textbox(
        slide,
        left1 + Inches(0.15), content_top + Inches(0.48),
        panel_w - Inches(0.3), Inches(0.22),
        "File: ARRGT_Split1_Bivariate.jrn",
        font_size=8, bold=False, color=MID_GREY,
    )
    _add_textbox(
        slide,
        left1 + Inches(0.15), content_top + Inches(0.72),
        panel_w - Inches(0.3), Inches(0.22),
        "Double-click the icon below to open in JMP",
        font_size=8, bold=False, color=MID_GREY,
    )
    if img1 and img1.exists():
        slide.shapes.add_picture(str(img1),
                                 left1 + Inches(0.15), img_top, img_w, img_h)

    # ── Split2 Bivariate panel ───────────────────────────────────────────────
    _add_rect(slide, left2, content_top, panel_w, content_h, LIGHT_BLUE)
    _add_textbox(
        slide,
        left2 + Inches(0.15), content_top + Inches(0.12),
        panel_w - Inches(0.3), Inches(0.30),
        "ARRGT Split2 Bivariate Analysis",
        font_size=10, bold=True, color=INTEL_BLUE,
    )
    _add_textbox(
        slide,
        left2 + Inches(0.15), content_top + Inches(0.48),
        panel_w - Inches(0.3), Inches(0.22),
        "File: ARRGT_Split2_Bivariate.jrn",
        font_size=8, bold=False, color=MID_GREY,
    )
    _add_textbox(
        slide,
        left2 + Inches(0.15), content_top + Inches(0.72),
        panel_w - Inches(0.3), Inches(0.22),
        "Double-click the icon below to open in JMP",
        font_size=8, bold=False, color=MID_GREY,
    )
    if img2 and img2.exists():
        slide.shapes.add_picture(str(img2),
                                 left2 + Inches(0.15), img_top, img_w, img_h)


# ---------------------------------------------------------------------------
# Page 5 – MIO Results
# ---------------------------------------------------------------------------
def build_MIO_result(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ── Header ───────────────────────────────────────────────────────────────
    _add_rect(slide, 0, 0, SLIDE_W, Inches(1.0), INTEL_BLUE)
    _add_textbox(
        slide,
        Inches(0.3), Inches(0.15),
        Inches(12), Inches(0.7),
        "Brita Yield Report  |  MIO Results",
        font_size=24, bold=True, color=WHITE,
    )

    # ── Footer ───────────────────────────────────────────────────────────────
    _add_rect(slide, 0, SLIDE_H - Inches(0.3), SLIDE_W, Inches(0.3), INTEL_BLUE)
    _add_textbox(
        slide,
        Inches(0.3), SLIDE_H - Inches(0.28),
        Inches(12), Inches(0.25),
        "Intel Confidential",
        font_size=8, color=WHITE,
    )

    content_top = Inches(1.08)
    content_h   = SLIDE_H - content_top - Inches(0.38)
    right_w     = Inches(2.8)
    right_left  = SLIDE_W - right_w - Inches(0.2)
    tbl_left    = Inches(0.3)
    tbl_w       = right_left - tbl_left - Inches(0.2)

    # ── Filtered yield table ─────────────────────────────────────────────────
    df = pd.read_csv(MIO_FILTERED_CSV, dtype=str).fillna("")
    cols = list(df.columns)
    n_data = len(df)
    n_rows = n_data + 1   # header + data

    col_widths = {
        "Material":    Inches(1.4),
        "TestProgram": Inches(1.5),
        "Name":        Inches(2.8),
        "Module":      Inches(1.5),
        "Total":       Inches(0.6),
        "#P":          Inches(0.55),
        "#F":          Inches(0.55),
        "#U":          Inches(0.55),
        "%P":          Inches(0.65),
    }
    n_cols = len(cols)
    row_h  = Inches(0.28)
    tbl_h  = min(row_h * n_rows, content_h)

    table = _add_table(slide, n_rows, n_cols,
                       tbl_left, content_top, tbl_w, tbl_h)

    # Set column widths
    remaining = tbl_w
    for ci, col in enumerate(cols):
        w = col_widths.get(col, Inches(1.0))
        if ci < n_cols - 1:
            table.columns[ci].width = w
            remaining -= w
        else:
            table.columns[ci].width = max(remaining, Inches(0.5))

    # Header row
    for ci, col in enumerate(cols):
        _style_cell(table.cell(0, ci), col,
                    font_size=9, bold=True,
                    fg=WHITE, bg=INTEL_BLUE, align=PP_ALIGN.CENTER)

    # Data rows
    for ri, (_, row) in enumerate(df.iterrows()):
        bg = LIGHT_BLUE if ri % 2 == 0 else WHITE
        for ci, col in enumerate(cols):
            val = str(row[col])
            align = PP_ALIGN.CENTER if col in ("Total", "#P", "#F", "#U", "%P") else PP_ALIGN.LEFT
            _style_cell(table.cell(ri + 1, ci), val,
                        font_size=8, bold=False,
                        fg=DARK_GREY, bg=bg, align=align)

    # ── Right panel: OLE attachment placeholder ───────────────────────────────
    _add_rect(slide, right_left, content_top, right_w, content_h, LIGHT_BLUE)
    _add_textbox(
        slide,
        right_left + Inches(0.15), content_top + Inches(0.12),
        right_w - Inches(0.3), Inches(0.30),
        "MIO Full Yield Data",
        font_size=10, bold=True, color=INTEL_BLUE,
    )
    _add_textbox(
        slide,
        right_left + Inches(0.15), content_top + Inches(0.48),
        right_w - Inches(0.3), Inches(0.22),
        "File: MIO_yield.csv",
        font_size=8, bold=False, color=MID_GREY,
    )
    _add_textbox(
        slide,
        right_left + Inches(0.15), content_top + Inches(0.72),
        right_w - Inches(0.3), Inches(0.22),
        "Double-click the icon below to open the data",
        font_size=8, bold=False, color=MID_GREY,
    )


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------
def main():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    cfg          = _load_config()
    check_period = cfg.get("check_period", 10)
    check_count  = cfg.get("check_count", None)

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print("[page 1] Checking required files...")
    _wait_for_files([XEUS_CSV, LOT_WAFER_RESULT], check_period, check_count)
    build_test_run_summary(prs)

    print("[page 2] Checking required files...")
    _wait_for_files([SOC_SCAN_HRY_CSV, SOC_SCAN_HRY_XLSX], check_period, check_count)
    build_SOC_SCAN_result(prs)

    print("[page 3] Checking required files...")
    _wait_for_files([MBIST_AGG_CSV, MBIST_FILTERED_CSV], check_period, check_count)
    build_MBIST_result(prs)

    print("[page 4] Checking required files...")
    _wait_for_files([ARRGT_SPLIT1_JRN, ARRGT_SPLIT2_JRN,
                     ARRGT_SPLIT1_PPTX, ARRGT_SPLIT2_PPTX], check_period, check_count)
    print("[page 4] Exporting slide images from ARRGT PPTXes...")
    _export_slide_as_image(ARRGT_SPLIT1_PPTX, ARRGT_SPLIT1_PNG)
    _export_slide_as_image(ARRGT_SPLIT2_PPTX, ARRGT_SPLIT2_PNG)
    build_ARRGT_result(prs,
                       img1=ARRGT_SPLIT1_PNG if ARRGT_SPLIT1_PNG.exists() else None,
                       img2=ARRGT_SPLIT2_PNG if ARRGT_SPLIT2_PNG.exists() else None)

    print("[page 5] Checking required files...")
    _wait_for_files([MIO_FILTERED_CSV, MIO_YIELD_CSV], check_period, check_count)
    build_MIO_result(prs)

    prs.save(PPT_OUT)
    print(f"Report saved -> {PPT_OUT}")

    # Embed OLE attachments (COM requires absolute paths internally;
    # path constants above are all relative to the working directory)
    _embed_csv_ole(PPT_OUT, SOC_SCAN_HRY_XLSX,   slide_index=1)
    _embed_csv_ole(PPT_OUT, MBIST_FILTERED_CSV,  slide_index=2)
    _embed_csv_ole(PPT_OUT, ARRGT_SPLIT1_JRN, slide_index=3,
                   left_in=0.35, top_in=5.5, width_in=2.2, height_in=1.6)
    _embed_csv_ole(PPT_OUT, ARRGT_SPLIT2_JRN, slide_index=3,
                   left_in=6.95, top_in=5.5, width_in=2.2, height_in=1.6)
    _embed_csv_ole(PPT_OUT, MIO_YIELD_CSV, slide_index=4,
                   left_in=10.55, top_in=1.95, width_in=2.2, height_in=1.6)


if __name__ == "__main__":
    main()
